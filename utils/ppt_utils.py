from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO

class PptUtils:

    def __init__(self,slide):
        self.slide = slide


    def cm_to_inches(self, cm_value):
        return cm_value / 2.54

    def cm_to_emu(self, cm_value):
        inches_value = self.cm_to_inches(cm_value)
        return int(inches_value * 914400)
    
    def add_textbox_to_slide(self, text, left, top, width, height, font_family,font_size, font_color=(0, 0, 0),alignment='center'):
        
        textbox = self.slide.shapes.add_textbox(left=self.cm_to_emu(left), 
                                                top=self.cm_to_emu(top), 
                                                width=self.cm_to_emu(width), 
                                                height=self.cm_to_emu(height))
        
        frame = textbox.text_frame
        frame.fit_text(font_family=font_family, max_size=self.cm_to_emu(100))
        frame.word_wrap=True
        p = frame.paragraphs[0]
        if alignment=='center':
            p.alignment = PP_ALIGN.CENTER
        parts = text.split('*')
        print(parts)
        for i, part in enumerate(parts):
            if len(part)<2:
                continue
            else:
                run = p.add_run()
                run.text = part
                font = run.font
                font.name = font_family
                font.size = Pt(font_size)
                font.color.rgb = RGBColor(*font_color)
                if i % 2 == 1:
                    font.bold = True

    def add_image_to_slide(self,image, left, top, width, height):
        img_stream = BytesIO()
        image.save(img_stream,format='PNG')
         # Rewind the stream to the beginning
        img_stream.seek(0)
        # Add the image from BytesIO to the slide
        self.slide.shapes.add_picture(img_stream,left=self.cm_to_emu(left), 
                                            top=self.cm_to_emu(top), 
                                            width=self.cm_to_emu(width), 
                                            height=self.cm_to_emu(height))
